import streamlit as st
import PyPDF2
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from nltk.tokenize import word_tokenize

def read_text(file):
    return file.read().decode('utf-8')

def read_docx(file):
    docx = Document(file)
    text = ""
    for paragraph in docx.paragraphs:
        text += paragraph.text
    return text

def read_pdf(uploaded_file):
    text = ""
    with uploaded_file as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    return text

def read_pptx(uploaded_file):
    text = ""
    prs = Presentation(uploaded_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def calculate_cosine_similarity(text1, text2):
    documents = [text1, text2]
    tfidf_vectorizer = TfidfVectorizer()
    tfidf_matrix = tfidf_vectorizer.fit_transform(documents)
    similarity_matrix = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix)
    return similarity_matrix[0][1]

def calculate_jaccard_similarity(text1, text2):
    set1 = set(word_tokenize(text1.lower()))
    set2 = set(word_tokenize(text2.lower()))
    intersection = len(set1.intersection(set2))
    union = len(set1.union(set2))
    return intersection / union

def calculate_overlap_coefficient(text1, text2):
    set1 = set(word_tokenize(text1.lower()))
    set2 = set(word_tokenize(text2.lower()))
    intersection = len(set1.intersection(set2))
    min_len = min(len(set1), len(set2))
    return intersection / min_len

def main():
    st.title("Document Similarity Calculator")

    uploaded_file1 = st.file_uploader("Upload first document", type=["txt", "docx", "pdf", "pptx"])
    uploaded_file2 = st.file_uploader("Upload second document", type=["txt", "docx", "pdf", "pptx"])
    threshold = st.slider("Set threshold", 0.0, 1.0, 0.5)

    if uploaded_file1 is not None and uploaded_file2 is not None:
        text1, text2 = "", ""
        if uploaded_file1.type == "text/plain":
            text1 = read_text(uploaded_file1)
        elif uploaded_file1.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text1 = read_docx(uploaded_file1)
        elif uploaded_file1.type == "application/pdf":
            text1 = read_pdf(uploaded_file1)
        elif uploaded_file1.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text1 = read_pptx(uploaded_file1)

        if uploaded_file2.type == "text/plain":
            text2 = read_text(uploaded_file2)
        elif uploaded_file2.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text2 = read_docx(uploaded_file2)
        elif uploaded_file2.type == "application/pdf":
            text2 = read_pdf(uploaded_file2)
        elif uploaded_file2.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text2 = read_pptx(uploaded_file2)

        cosine_sim = calculate_cosine_similarity(text1, text2)
        jaccard_sim = calculate_jaccard_similarity(text1, text2)
        overlap_coeff = calculate_overlap_coefficient(text1, text2)

        st.subheader("Similarity Results:")
        
        if cosine_sim >= threshold:
            st.write(f"Cosine Similarity: *{cosine_sim:.2f}*", unsafe_allow_html=True)
        else:
            st.write(f"Cosine Similarity: <font color='red' size='+2'>{cosine_sim:.2f}</font>", unsafe_allow_html=True)

        if jaccard_sim >= threshold:
            st.write(f"Jaccard Similarity: *{jaccard_sim:.2f}*", unsafe_allow_html=True)
        else:
            st.write(f"Jaccard Similarity: <font color='red' size='+2'>{jaccard_sim:.2f}</font>", unsafe_allow_html=True)

        if overlap_coeff >= threshold:
            st.write(f"Overlap Coefficient: *{overlap_coeff:.2f}*", unsafe_allow_html=True)
        else:
            st.write(f"Overlap Coefficient: <font color='red' size='+2'>{overlap_coeff:.2f}</font>", unsafe_allow_html=True)

        if cosine_sim >= threshold or jaccard_sim >= threshold or overlap_coeff >= threshold:
            st.markdown("<p style='color:red; font-size:30px;'>The documents are similar.</p>", unsafe_allow_html=True)
        else:
            st.markdown("<p style='color:green; font-size:30px;'>The documents are not similar.</p>", unsafe_allow_html=True)

if __name__ == "__main__":
    main()